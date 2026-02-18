"""
PowerPoint Presentation Generator
X-Health Default Prediction Model - Kognita Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # slide backgrounds / headers
TEAL        = RGBColor(0x00, 0xB4, 0xD8)   # accent / titles
LIGHT_TEAL  = RGBColor(0x90, 0xE0, 0xEF)   # sub-headings / bullet accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # content area bg
DARK_GRAY   = RGBColor(0x2D, 0x3A, 0x4A)   # body text on light bg
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)   # highlight / warning
GREEN       = RGBColor(0x06, 0xD6, 0x7E)   # success / positive metric
MID_NAVY    = RGBColor(0x1A, 0x2E, 0x44)   # card bg on dark slides


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=14, color=WHITE, accent=TEAL,
                   font_name="Calibri", line_spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet character
        run0 = p.add_run()
        run0.text = "▶  "
        run0.font.size = Pt(font_size - 1)
        run0.font.color.rgb = accent
        run0.font.name = font_name
        # body text
        run1 = p.add_run()
        run1.text = item
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = color
        run1.font.name = font_name
    return txb


def add_metric_card(slide, label, value, l, t, w=1.8, h=1.1,
                    bg=MID_NAVY, value_color=TEAL, label_color=LIGHT_TEAL):
    add_rect(slide, l, t, w, h, bg)
    add_text_box(slide, value, l + 0.05, t + 0.05, w - 0.1, 0.55,
                 font_size=26, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l + 0.05, t + 0.60, w - 0.1, 0.45,
                 font_size=11, bold=False, color=label_color, align=PP_ALIGN.CENTER)


def slide_header(slide, title, subtitle=None, title_y=0.18, title_fs=28, sub_fs=15):
    add_text_box(slide, title, 0.5, title_y, 9.0, 0.55,
                 font_size=title_fs, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, title_y + 0.50, 9.0, 0.35,
                     font_size=sub_fs, bold=False, color=LIGHT_TEAL, align=PP_ALIGN.LEFT)
    # accent underline
    add_rect(slide, 0.5, title_y + (0.88 if subtitle else 0.62), 9.0, 0.04, TEAL)


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_NAVY)

    # Full-width top accent bar
    add_rect(slide, 0, 0, 10, 0.12, TEAL)

    # Left colour block
    add_rect(slide, 0, 0.12, 0.18, 7.38, TEAL)

    # Main title
    add_text_box(slide,
                 "X-Health Default\nPrediction Model",
                 0.35, 1.4, 9.3, 1.6,
                 font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Sub-title
    add_text_box(slide,
                 "Machine Learning for Credit Risk in B2B Healthcare Equipment Sales",
                 0.35, 3.15, 9.3, 0.5,
                 font_size=18, bold=False, color=LIGHT_TEAL, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(slide, 0.35, 3.75, 9.3, 0.04, TEAL)

    # Meta info row
    meta = [
        ("Client",   "X-Health"),
        ("Team",     "Kognita Lab"),
        ("Date",     "February 2026"),
        ("Version",  "v3.4.0"),
    ]
    for i, (lbl, val) in enumerate(meta):
        x = 0.35 + i * 2.4
        add_text_box(slide, lbl.upper(), x, 3.90, 2.2, 0.25,
                     font_size=9, bold=True, color=LIGHT_TEAL)
        add_text_box(slide, val, x, 4.15, 2.2, 0.35,
                     font_size=14, bold=True, color=WHITE)

    # Bottom bar
    add_rect(slide, 0, 7.38, 10, 0.12, TEAL)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Agenda", "What we will cover today")

    topics = [
        ("01", "Business Context & Problem Statement"),
        ("02", "Dataset Overview"),
        ("03", "Exploratory Data Analysis"),
        ("04", "Feature Engineering"),
        ("05", "Modeling & Evaluation"),
        ("06", "Model Interpretability (SHAP)"),
        ("07", "Threshold Optimization"),
        ("08", "Results & Key Metrics"),
        ("09", "Prediction Pipeline"),
        ("10", "Conclusions & Next Steps"),
    ]

    col_w = 4.5
    for i, (num, topic) in enumerate(topics):
        col = i % 2
        row = i // 2
        x = 0.5 + col * col_w
        y = 1.55 + row * 1.1

        add_rect(slide, x, y, col_w - 0.2, 0.85, MID_NAVY)
        add_text_box(slide, num, x + 0.1, y + 0.05, 0.5, 0.75,
                     font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text_box(slide, topic, x + 0.65, y + 0.18, col_w - 0.95, 0.55,
                     font_size=13, bold=False, color=WHITE)


def slide_business_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Business Context", "01 — Problem Statement")

    # Left panel
    add_rect(slide, 0.5, 1.35, 4.4, 5.8, MID_NAVY)
    add_text_box(slide, "The Challenge", 0.65, 1.45, 4.1, 0.4,
                 font_size=16, bold=True, color=TEAL)
    add_bullet_box(slide, [
        "X-Health sells health electronic devices in B2B credit-based transactions",
        "Buyers pay in lump sum or installments after delivery",
        "A significant share of transactions end in payment default (calote)",
        "Defaults directly impact revenue, cash flow, and financial planning",
        "Manual review is costly and inconsistent at scale",
    ], 0.65, 1.90, 4.1, 4.5, font_size=13, color=WHITE)

    # Right panel
    add_rect(slide, 5.1, 1.35, 4.4, 5.8, MID_NAVY)
    add_text_box(slide, "Our Solution", 5.25, 1.45, 4.1, 0.4,
                 font_size=16, bold=True, color=GREEN)
    add_bullet_box(slide, [
        "Build a probabilistic ML model to predict default risk",
        "Assess each purchase order BEFORE finalisation",
        "Output: probability score + binary decision",
        "Support 3 threshold modes for different risk strategies",
        "Enable data-driven credit approval at scale",
    ], 5.25, 1.90, 4.1, 4.5, font_size=13, color=WHITE, accent=GREEN)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Dataset Overview", "02 — Data at a Glance")

    # Top metric cards
    metrics = [
        ("122,500+", "Transactions"),
        ("2017–2019", "Time Period"),
        ("22", "Raw Features"),
        ("~16–17%", "Default Rate"),
        ("22 MB", "Dataset Size"),
    ]
    for i, (val, lbl) in enumerate(metrics):
        add_metric_card(slide, lbl, val, 0.2 + i * 1.96, 1.35, w=1.82, h=1.1)

    # Feature categories
    cats = [
        ("Internal (Client History)", TEAL, [
            "default_3months — Past defaults",
            "ioi_36months / ioi_3months — Order frequency",
            "valor_por_vencer — Pending payments",
            "valor_vencido — Overdue amount",
            "valor_quitado — Historically paid",
        ]),
        ("External (Credit Bureau)", LIGHT_TEAL, [
            "quant_protestos — # of legal protests",
            "valor_protestos — Protest amount (R$)",
            "6 variables excluded (≥99% zeros)",
            "  ▷ participacao_falencia_valor (100%)",
            "  ▷ falencia_concordata_qtd (99.95%)",
        ]),
        ("Categorical / Context", ORANGE, [
            "tipo_sociedade — Business type",
            "opcao_tributaria — Tax regime",
            "atividade_principal — Main activity",
            "forma_pagamento — Payment form",
            "valor_total_pedido, month",
        ]),
    ]
    col_w = 3.1
    for i, (title, col, items) in enumerate(cats):
        x = 0.25 + i * 3.25
        add_rect(slide, x, 2.65, col_w, 4.45, MID_NAVY)
        add_text_box(slide, title, x + 0.1, 2.72, col_w - 0.15, 0.38,
                     font_size=13, bold=True, color=col)
        add_rect(slide, x + 0.1, 3.12, col_w - 0.2, 0.03, col)
        add_bullet_box(slide, items, x + 0.1, 3.20, col_w - 0.15, 3.8,
                       font_size=11, color=WHITE, accent=col)


def slide_eda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Exploratory Data Analysis", "03 — Key Findings")

    findings = [
        (TEAL,   "Maturity Bias Detected",
         "Recent months show artificially low default rates (8–10%) vs mature cohorts (18–20%) because payment cycles are incomplete. StratifiedKFold CV was chosen over out-of-time split to avoid penalising the model unfairly."),
        (GREEN,  "Statistical Significance",
         "Mann-Whitney U tests with effect sizes confirmed that overdue amounts, protest counts, and order frequency are statistically significant predictors of default."),
        (ORANGE, "Data Quality Issues",
         "144 records (0.12%) with negative order values removed — likely data-entry errors or cancellations. 6 Serasa variables with ≥99% zeros excluded as they carry no predictive power."),
        (LIGHT_TEAL, "Cluster Profiles (K-Means + PCA)",
         "K-Means clustering revealed distinct risk groups: low-risk frequent buyers, high-risk first-timers with overdue balances, and medium-risk clients with legal protests."),
    ]

    for i, (col, title, desc) in enumerate(findings):
        row = i // 2
        c   = i % 2
        x = 0.4 + c * 4.8
        y = 1.35 + row * 2.85
        add_rect(slide, x, y, 4.55, 2.65, MID_NAVY)
        add_rect(slide, x, y, 0.07, 2.65, col)
        add_text_box(slide, title, x + 0.2, y + 0.12, 4.25, 0.4,
                     font_size=15, bold=True, color=col)
        add_text_box(slide, desc, x + 0.2, y + 0.58, 4.25, 1.95,
                     font_size=12, bold=False, color=WHITE)


def slide_feature_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Feature Engineering", "04 — 7 Business-Driven Derived Features")

    features = [
        ("total_exposto",        "Total financial exposure",      "vencido + por_vencer + quitado",            TEAL),
        ("razao_inadimplencia",  "Delinquency ratio",             "vencido / (total_exposto + 1)",             ORANGE),
        ("taxa_cobertura_divida","Historical payment capacity",   "quitado / (pedido + 1)",                    GREEN),
        ("razao_vencido_pedido", "Current leverage ratio",        "vencido / (pedido + 1)",                    LIGHT_TEAL),
        ("flag_risco_juridico",  "Binary legal-risk flag",        "1 if quant_protestos > 0 else 0",           ORANGE),
        ("ticket_medio_protestos","Avg. protest severity",        "valor_protestos / (quant_protestos + 1)",   TEAL),
        ("qtd_parcelas",         "Payment-complexity proxy",      "count('/') in forma_pagamento + 1",         GREEN),
    ]

    col_w = 4.55
    for i, (name, label, formula, col) in enumerate(features):
        row = i // 2
        c   = i % 2
        x = 0.25 + c * 4.8
        y = 1.35 + row * 1.55

        add_rect(slide, x, y, col_w, 1.35, MID_NAVY)
        add_rect(slide, x, y, col_w, 0.06, col)
        add_text_box(slide, name, x + 0.12, y + 0.10, col_w - 0.2, 0.32,
                     font_size=13, bold=True, color=col)
        add_text_box(slide, label, x + 0.12, y + 0.42, col_w - 0.2, 0.28,
                     font_size=11, bold=False, color=LIGHT_TEAL)
        add_text_box(slide, f"Formula: {formula}", x + 0.12, y + 0.72, col_w - 0.2, 0.55,
                     font_size=10, bold=False, color=WHITE)

    # 7th feature (odd one out) centred at bottom
    if len(features) % 2 == 1:
        pass  # already handled — index 6 lands at row=3, col=0 (left side)


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Preprocessing Pipeline", "04b — ColumnTransformer Architecture")

    steps = [
        ("Raw Input\n18 features", DARK_GRAY, WHITE),
        ("ColumnTransformer", MID_NAVY, TEAL),
        ("Processed Output\n28 features", MID_NAVY, GREEN),
        ("XGBoost\nClassifier", MID_NAVY, ORANGE),
        ("Default\nProbability", MID_NAVY, TEAL),
    ]

    box_w, box_h, gap = 1.55, 1.0, 0.25
    total = len(steps) * box_w + (len(steps) - 1) * gap
    start_x = (10 - total) / 2
    y = 1.65

    for i, (label, bg, fc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, y, box_w, box_h, bg, TEAL if i != 0 else DARK_GRAY, 1)
        add_text_box(slide, label, x + 0.05, y + 0.1, box_w - 0.1, box_h - 0.15,
                     font_size=12, bold=True, color=fc, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + 0.02
            add_text_box(slide, "→", ax, y + 0.28, gap, 0.5,
                         font_size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Three column transformer sub-steps
    sub_cols = [
        ("Numeric (16 features)", TEAL, [
            "SimpleImputer (median)",
            "RobustScaler",
        ]),
        ("High-Cardinality (3)", ORANGE, [
            "SimpleImputer",
            "Target Encoding",
            "tipo_sociedade",
            "atividade_principal",
            "forma_pagamento",
        ]),
        ("Low-Cardinality (2)", GREEN, [
            "SimpleImputer",
            "OneHotEncoder",
            "opcao_tributaria",
            "month",
        ]),
    ]
    for i, (title, col, items) in enumerate(sub_cols):
        x = 1.0 + i * 2.95
        y2 = 3.05
        add_rect(slide, x, y2, 2.8, 3.9, MID_NAVY)
        add_rect(slide, x, y2, 2.8, 0.06, col)
        add_text_box(slide, title, x + 0.1, y2 + 0.1, 2.6, 0.35,
                     font_size=12, bold=True, color=col)
        add_bullet_box(slide, items, x + 0.1, y2 + 0.5, 2.6, 3.2,
                       font_size=11, color=WHITE, accent=col)


def slide_modeling(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Model Comparison", "05 — 6 Algorithms Evaluated with StratifiedKFold CV")

    models = [
        ("Logistic Regression", "0.82", "Baseline linear model"),
        ("Decision Tree",       "0.84", "Interpretable but overfit"),
        ("Random Forest",       "0.90", "Solid ensemble performance"),
        ("Gradient Boosting",   "0.91", "Strong generalisation"),
        ("XGBoost",             "0.92", "★ Best overall — Selected", True),
        ("LightGBM",            "0.91", "Close runner-up"),
    ]

    col_w = 4.45
    for i, row_data in enumerate(models):
        name, auc, note = row_data[0], row_data[1], row_data[2]
        winner = len(row_data) == 4
        c   = i % 2
        r   = i // 2
        x = 0.3 + c * 4.8
        y = 1.35 + r * 1.85

        bg  = RGBColor(0x06, 0x36, 0x06) if winner else MID_NAVY
        brd = GREEN if winner else None
        add_rect(slide, x, y, col_w, 1.65, bg, brd, 2 if winner else 0)

        fc = GREEN if winner else TEAL
        add_text_box(slide, name, x + 0.15, y + 0.12, col_w - 0.25, 0.38,
                     font_size=15, bold=True, color=fc)
        add_text_box(slide, f"ROC-AUC: {auc}", x + 0.15, y + 0.55, col_w - 0.25, 0.35,
                     font_size=20, bold=True, color=WHITE if not winner else GREEN)
        add_text_box(slide, note, x + 0.15, y + 0.96, col_w - 0.25, 0.55,
                     font_size=11, bold=False, color=LIGHT_TEAL if not winner else GREEN)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Model Performance", "08 — XGBoost Final Results")

    # Big metric cards
    top_metrics = [
        ("0.9251", "Test ROC-AUC", GREEN),
        ("0.9197", "CV ROC-AUC Mean", TEAL),
        ("±0.0021", "CV ROC-AUC Std", LIGHT_TEAL),
    ]
    for i, (val, lbl, col) in enumerate(top_metrics):
        x = 0.5 + i * 3.15
        add_rect(slide, x, 1.35, 2.95, 1.4, MID_NAVY)
        add_text_box(slide, val, x + 0.1, 1.42, 2.75, 0.7,
                     font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x + 0.1, 2.15, 2.75, 0.45,
                     font_size=13, bold=False, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

    # Secondary metrics
    sec = [
        ("F1-Score",  "0.6709"),
        ("Precision", "0.5691"),
        ("Recall",    "0.8170"),
        ("Threshold", "0.647"),
    ]
    for i, (lbl, val) in enumerate(sec):
        x = 0.5 + i * 2.35
        add_metric_card(slide, lbl, val, x, 2.95, w=2.15, h=1.05)

    # Interpretation bullets
    add_rect(slide, 0.5, 4.2, 9.1, 2.95, MID_NAVY)
    add_text_box(slide, "Interpretation", 0.65, 4.28, 8.8, 0.38,
                 font_size=15, bold=True, color=TEAL)
    add_bullet_box(slide, [
        "Excellent discrimination ability — ROC-AUC > 0.92 places this model in the top tier of credit risk systems",
        "Strong recall (81.7%) — captures most actual defaults, minimising financial exposure",
        "Acceptable precision (56.9%) — some false alarms, manageable with threshold tuning",
        "Highly stable across folds — CV std of ±0.21% confirms the model generalises reliably",
        "Class imbalance (~16% defaults) handled via class_weight='balanced' in XGBoost",
    ], 0.65, 4.70, 8.8, 2.35, font_size=12, color=WHITE)


def slide_shap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Model Interpretability — SHAP", "06 — Why Does the Model Predict Default?")

    # Top features ranking
    features_ranked = [
        ("razao_inadimplencia",     "Delinquency Ratio",          95, ORANGE),
        ("valor_vencido",           "Overdue Amount (R$)",         88, ORANGE),
        ("quant_protestos",         "# Legal Protests",            76, TEAL),
        ("razao_vencido_pedido",    "Overdue-to-Order Ratio",      70, TEAL),
        ("ioi_3months",             "Recency of Orders (3M)",      62, GREEN),
        ("default_3months",         "Past Defaults (3M)",          58, ORANGE),
        ("total_exposto",           "Total Financial Exposure",    50, TEAL),
        ("taxa_cobertura_divida",   "Historical Payment Capacity", 42, GREEN),
    ]

    add_text_box(slide, "Top Predictive Features (SHAP Importance)", 0.5, 1.35, 9.0, 0.38,
                 font_size=14, bold=True, color=TEAL)

    bar_x, bar_w_max = 3.4, 5.8
    for i, (fname, label, score, col) in enumerate(features_ranked):
        y = 1.85 + i * 0.68
        add_text_box(slide, label, 0.5, y + 0.05, 2.8, 0.45,
                     font_size=11, bold=False, color=WHITE, align=PP_ALIGN.RIGHT)
        bar_len = bar_w_max * score / 100
        add_rect(slide, bar_x, y + 0.08, bar_len, 0.38, col)
        add_text_box(slide, f"{score}", bar_x + bar_len + 0.1, y + 0.08, 0.6, 0.38,
                     font_size=11, bold=True, color=col)


def slide_thresholds(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Threshold Optimisation", "07 — 3 Configurable Decision Strategies")

    thresholds = [
        ("F1 — Balanced",
         "0.647",
         "Default Recommendation",
         "Balances precision and recall equally. Best for general-purpose deployment when both false positives (missed sales) and false negatives (uncollected debt) carry similar costs.",
         TEAL,
         ["Precision: 56.9%", "Recall: 81.7%", "F1-Score: 67.1%"]),
        ("F2 — Capture Defaults",
         "0.462",
         "Maximise Default Detection",
         "Prioritises recall over precision. Recommended when the financial cost of a missed default far exceeds the cost of rejecting a legitimate buyer.",
         ORANGE,
         ["Higher Recall", "Lower Precision", "Minimises Financial Loss"]),
        ("F0.5 — Avoid False Alarms",
         "0.793",
         "Minimise False Rejections",
         "Prioritises precision over recall. Best when retaining customers is critical and the business can tolerate occasional defaults.",
         GREEN,
         ["Higher Precision", "Lower Recall", "Maximises Sales Approval"]),
    ]

    for i, (title, thresh, subtitle, desc, col, stats) in enumerate(thresholds):
        x = 0.28 + i * 3.24
        add_rect(slide, x, 1.35, 3.1, 5.75, MID_NAVY)
        add_rect(slide, x, 1.35, 3.1, 0.07, col)
        add_text_box(slide, title, x + 0.12, 1.45, 2.85, 0.38,
                     font_size=14, bold=True, color=col)
        add_text_box(slide, f"Threshold: {thresh}", x + 0.12, 1.88, 2.85, 0.48,
                     font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, subtitle, x + 0.12, 2.42, 2.85, 0.35,
                     font_size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_rect(slide, x + 0.12, 2.82, 2.85, 0.03, col)
        add_text_box(slide, desc, x + 0.12, 2.90, 2.85, 1.65,
                     font_size=11, bold=False, color=WHITE)
        add_bullet_box(slide, stats, x + 0.12, 4.62, 2.85, 2.25,
                       font_size=12, color=WHITE, accent=col)


def slide_prediction_api(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Prediction Pipeline", "09 — How to Use the Model")

    # Left: input
    add_rect(slide, 0.3, 1.35, 4.5, 5.75, MID_NAVY)
    add_text_box(slide, "Input Features (dict)", 0.45, 1.42, 4.2, 0.38,
                 font_size=13, bold=True, color=TEAL)
    code_input = (
        'input_dict = {\n'
        '  "default_3months": 0,\n'
        '  "ioi_36months": 18.2,\n'
        '  "ioi_3months": 5.0,\n'
        '  "valor_por_vencer": 10000,\n'
        '  "valor_vencido": 50000,\n'
        '  "valor_quitado": 500000,\n'
        '  "quant_protestos": 0,\n'
        '  "valor_protestos": 0,\n'
        '  "valor_total_pedido": 35000,\n'
        '  "tipo_sociedade": "ltda",\n'
        '  "opcao_tributaria": "simples",\n'
        '  "atividade_principal": "...",\n'
        '  "forma_pagamento": "30/60/90",\n'
        '  "month": 6\n'
        '}'
    )
    add_text_box(slide, code_input, 0.45, 1.88, 4.2, 5.0,
                 font_size=10, bold=False, color=LIGHT_TEAL, font_name="Courier New")

    # Right: output
    add_rect(slide, 5.0, 1.35, 4.6, 5.75, MID_NAVY)
    add_text_box(slide, "Output & Function Call", 5.15, 1.42, 4.3, 0.38,
                 font_size=13, bold=True, color=GREEN)
    add_text_box(slide,
                 "result = predict_default(\n    input_dict,\n    threshold=0.647\n)",
                 5.15, 1.88, 4.3, 0.85,
                 font_size=11, bold=False, color=LIGHT_TEAL, font_name="Courier New")
    add_text_box(slide,
                 '# Returns:\n{\n  "default": 0 or 1,\n  "probability": 0.0–1.0,\n  "confidence": 0.0–1.0,\n  "threshold_used": 0.647\n}',
                 5.15, 2.82, 4.3, 1.5,
                 font_size=11, bold=False, color=LIGHT_TEAL, font_name="Courier New")

    add_text_box(slide, "Automatic Feature Engineering", 5.15, 4.48, 4.3, 0.35,
                 font_size=13, bold=True, color=TEAL)
    add_bullet_box(slide, [
        "Derives all 7 engineered features automatically",
        "Handles missing values & unknown categories",
        "Validates negative order values",
        "Supports batch prediction via predict_default_batch()",
    ], 5.15, 4.90, 4.3, 2.05, font_size=11, color=WHITE)


def slide_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    slide_header(slide, "Conclusions & Next Steps", "10 — Summary & Roadmap")

    # Achievements
    add_rect(slide, 0.3, 1.35, 5.9, 3.6, MID_NAVY)
    add_text_box(slide, "✓  Key Achievements", 0.45, 1.42, 5.6, 0.38,
                 font_size=14, bold=True, color=GREEN)
    add_bullet_box(slide, [
        "ROC-AUC of 0.9251 — excellent credit-risk discrimination",
        "Recall of 81.7% — captures most actual defaults",
        "Maturity bias identified and mitigated (StratifiedKFold)",
        "7 business-driven features engineered from raw data",
        "SHAP analysis provides full prediction explainability",
        "3 threshold modes for flexible business deployment",
        "Production-ready: serialised pipeline + standalone script",
    ], 0.45, 1.88, 5.6, 2.95, font_size=12, color=WHITE, accent=GREEN)

    # Roadmap
    add_rect(slide, 6.4, 1.35, 3.25, 3.6, MID_NAVY)
    add_text_box(slide, "→  Roadmap", 6.55, 1.42, 3.0, 0.38,
                 font_size=14, bold=True, color=ORANGE)
    add_bullet_box(slide, [
        "Probability calibration (Platt / Isotonic)",
        "Ensemble methods (stacking / blending)",
        "REST API via FastAPI + Docker",
        "Real-time monitoring & drift alerts",
        "Model Cards & fairness analysis",
        "Model Registry & versioning",
    ], 6.55, 1.88, 3.0, 2.95, font_size=12, color=WHITE, accent=ORANGE)

    # Tech stack row
    add_rect(slide, 0.3, 5.15, 9.35, 1.95, MID_NAVY)
    add_text_box(slide, "Technology Stack", 0.45, 5.22, 9.0, 0.38,
                 font_size=13, bold=True, color=TEAL)
    stack = ["Python 3.8+", "XGBoost 2.0", "scikit-learn 1.3", "SHAP 0.43",
             "LightGBM 4.0", "Pandas 2.0", "Matplotlib / Seaborn", "joblib"]
    for i, tech in enumerate(stack):
        x = 0.45 + i * 1.18
        add_rect(slide, x, 5.65, 1.1, 0.38, DARK_NAVY)
        add_text_box(slide, tech, x + 0.04, 5.67, 1.05, 0.34,
                     font_size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    add_rect(slide, 0, 0, 10, 0.12, TEAL)
    add_rect(slide, 0, 7.38, 10, 0.12, TEAL)
    add_rect(slide, 0, 0.12, 0.18, 7.26, TEAL)

    add_text_box(slide, "Thank You", 0.4, 2.2, 9.2, 1.0,
                 font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "X-Health Default Prediction Model  |  Kognita Lab  |  February 2026",
                 0.4, 3.35, 9.2, 0.45,
                 font_size=16, bold=False, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.0, 3.95, 6.0, 0.04, TEAL)
    add_text_box(slide, "Questions welcome",
                 0.4, 4.1, 9.2, 0.45,
                 font_size=18, bold=False, color=TEAL, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def build_presentation(output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_business_context(prs)
    slide_dataset(prs)
    slide_eda(prs)
    slide_feature_engineering(prs)
    slide_pipeline(prs)
    slide_modeling(prs)
    slide_metrics(prs)
    slide_shap(prs)
    slide_thresholds(prs)
    slide_prediction_api(prs)
    slide_conclusions(prs)
    slide_thank_you(prs)

    prs.save(output_path)
    print(f"✅  Presentation saved → {output_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    out = os.path.join(os.path.dirname(__file__), "..", "X-Health_Default_Prediction_Model.pptx")
    out = os.path.normpath(out)
    build_presentation(out)
